"""
Tests for Features 1–5:
  1 – DataFrame accessor (df.glyphx.*)
  2 – Type annotations + py.typed
  3 – CLI tool (glyphx plot / suggest / version)
  4 – PPTX export (fig.save("*.pptx"))
  5 – Accessibility (ARIA, alt text, keyboard nav, tabindex)
"""
from __future__ import annotations

import os
import re
import sys
import tempfile
from pathlib import Path
from io import StringIO

import numpy as np
import pandas as pd
import sys as _sys, os as _os
_sys.path.insert(0, _os.path.dirname(_os.path.abspath(__file__)))
try:
    import pytest
except ImportError:
    import pytest_shim as pytest  # noqa: F401
    _sys.modules["pytest"] = pytest

import glyphx
from glyphx import Figure
from glyphx.series import LineSeries, BarSeries, ScatterSeries


# ── Fixtures ─────────────────────────────────────────────────────────────────

@pytest.fixture
def sample_df():
    return pd.DataFrame({
        "month":   ["Jan", "Feb", "Mar", "Apr", "May", "Jun"],
        "revenue": [120, 135, 98, 170, 145, 190],
        "costs":   [80,  90,  60, 100,  95, 110],
        "region":  ["N", "S", "N", "S", "N", "S"],
    })

@pytest.fixture
def numeric_df():
    np.random.seed(42)
    return pd.DataFrame({
        "x": np.random.rand(50),
        "y": np.random.rand(50),
        "z": np.random.rand(50),
    })

@pytest.fixture
def basic_fig():
    fig = Figure(title="Test Chart", auto_display=False)
    fig.add(LineSeries([1, 2, 3], [4, 5, 6], label="Series A"))
    fig.add(BarSeries([1, 2, 3], [7, 3, 9], label="Series B"))
    return fig


# ============================================================
# Feature 1 — DataFrame Accessor
# ============================================================

class TestDataFrameAccessor:

    def test_accessor_registered(self, sample_df):
        assert hasattr(sample_df, "glyphx")

    def test_accessor_line(self, sample_df):
        fig = sample_df.glyphx.line(x="month", y="revenue", auto_display=False)
        assert isinstance(fig, Figure)
        assert "polyline" in fig.render_svg()

    def test_accessor_bar(self, sample_df):
        fig = sample_df.glyphx.bar(x="month", y="revenue", auto_display=False)
        assert "<rect" in fig.render_svg()

    def test_accessor_scatter(self, numeric_df):
        fig = numeric_df.glyphx.scatter(x="x", y="y", auto_display=False)
        assert "circle" in fig.render_svg()

    def test_accessor_hist(self, sample_df):
        fig = sample_df.glyphx.hist(col="revenue", bins=5, auto_display=False)
        assert "<rect" in fig.render_svg()

    def test_accessor_box(self, sample_df):
        fig = sample_df.glyphx.box(col="revenue", auto_display=False)
        assert "<svg" in fig.render_svg()

    def test_accessor_box_groupby(self, sample_df):
        fig = sample_df.glyphx.box(col="revenue", groupby="region", auto_display=False)
        assert "<svg" in fig.render_svg()

    def test_accessor_pie(self, sample_df):
        fig = sample_df.glyphx.pie(labels="region", values="revenue", auto_display=False)
        assert "<path" in fig.render_svg()

    def test_accessor_donut(self, sample_df):
        fig = sample_df.glyphx.donut(labels="region", values="revenue", auto_display=False)
        assert "<path" in fig.render_svg()

    def test_accessor_heatmap(self, numeric_df):
        small = numeric_df.head(5)
        fig = small.glyphx.heatmap(auto_display=False)
        assert "<rect" in fig.render_svg()

    def test_accessor_plot_dispatch_line(self, sample_df):
        fig = sample_df.glyphx.plot(kind="line", x="month", y="revenue", auto_display=False)
        assert "polyline" in fig.render_svg()

    def test_accessor_plot_dispatch_bar(self, sample_df):
        fig = sample_df.glyphx.plot(kind="bar", x="month", y="revenue", auto_display=False)
        assert "<rect" in fig.render_svg()

    def test_accessor_invalid_kind_raises(self, sample_df):
        with pytest.raises(ValueError, match="Unknown chart kind"):
            sample_df.glyphx.plot(kind="radar")

    def test_accessor_title_propagated(self, sample_df):
        fig = sample_df.glyphx.bar(
            x="month", y="revenue",
            title="Monthly Revenue", auto_display=False
        )
        assert "Monthly Revenue" in fig.render_svg()

    def test_accessor_theme_applied(self, sample_df):
        fig = sample_df.glyphx.line(
            x="month", y="revenue",
            theme="dark", auto_display=False
        )
        assert "#1e1e1e" in fig.render_svg()

    def test_accessor_xlabel_ylabel(self, sample_df):
        fig = sample_df.glyphx.line(
            x="month", y="revenue",
            xlabel="Month", ylabel="USD",
            auto_display=False
        )
        svg = fig.render_svg()
        assert "Month" in svg
        assert "USD" in svg

    def test_accessor_bar_groupby(self, sample_df):
        fig = sample_df.glyphx.bar(
            groupby="region", y="revenue", agg="sum",
            auto_display=False
        )
        assert "<rect" in fig.render_svg()

    def test_accessor_returns_figure_for_chaining(self, sample_df):
        """Accessor methods must return Figure so calls can be chained."""
        fig = (
            sample_df.glyphx
            .line(x="month", y="revenue", auto_display=False)
            .set_title("Chained Title")
        )
        assert isinstance(fig, Figure)
        assert fig.title == "Chained Title"

    def test_accessor_full_chain(self, sample_df, tmp_path):
        """End-to-end chain: accessor → set_theme → annotate → save."""
        out = str(tmp_path / "chain.svg")
        (
            sample_df.glyphx
            .bar(x="month", y="revenue", label="Revenue", auto_display=False)
            .set_theme("pastel")
            .set_title("Full Chain Test")
            .annotate("peak", x="Jun", y=190)
            .save(out)
        )
        assert os.path.exists(out)
        content = open(out).read()
        assert "Full Chain Test" in content


# ============================================================
# Feature 2 — Type Annotations + py.typed
# ============================================================

class TestTypeAnnotations:

    def test_py_typed_marker_exists(self):
        """PEP 561 requires a py.typed marker file in the package directory."""
        marker = Path(glyphx.__file__).parent / "py.typed"
        assert marker.exists(), "py.typed marker file missing"
        # File should be empty (or at most a brief comment)
        assert marker.stat().st_size < 512

    def test_figure_has_annotations(self):
        import inspect
        sig = inspect.signature(Figure.__init__)
        params = sig.parameters
        # Key parameters should have annotations
        assert params["width"].annotation != inspect.Parameter.empty
        assert params["height"].annotation != inspect.Parameter.empty
        assert params["title"].annotation != inspect.Parameter.empty
        assert params["theme"].annotation != inspect.Parameter.empty

    def test_figure_methods_have_return_annotations(self):
        import inspect
        for method_name in ("add", "show", "save", "share", "annotate",
                            "set_title", "set_theme", "set_size",
                            "set_xlabel", "set_ylabel", "set_legend"):
            method = getattr(Figure, method_name)
            sig    = inspect.signature(method)
            assert sig.return_annotation != inspect.Parameter.empty, \
                f"Figure.{method_name} missing return annotation"

    def test_add_returns_figure(self):
        fig = Figure(auto_display=False)
        result = fig.add(LineSeries([1, 2], [3, 4]))
        assert result is fig

    def test_annotate_returns_figure(self):
        fig = Figure(auto_display=False)
        fig.add(LineSeries([1, 2], [3, 4]))
        result = fig.annotate("label", x=1, y=3)
        assert result is fig

    def test_set_title_returns_figure(self):
        fig = Figure(auto_display=False)
        result = fig.set_title("New Title")
        assert result is fig
        assert fig.title == "New Title"

    def test_set_theme_returns_figure(self):
        fig = Figure(auto_display=False)
        result = fig.set_theme("dark")
        assert result is fig
        assert fig.theme["background"] == "#1e1e1e"

    def test_set_size_returns_figure(self):
        fig = Figure(auto_display=False)
        result = fig.set_size(1024, 768)
        assert result is fig
        assert fig.width == 1024
        assert fig.height == 768

    def test_set_xlabel_returns_figure(self):
        fig = Figure(auto_display=False)
        result = fig.set_xlabel("Time")
        assert result is fig
        assert fig.axes.xlabel == "Time"

    def test_set_ylabel_returns_figure(self):
        fig = Figure(auto_display=False)
        result = fig.set_ylabel("Revenue")
        assert result is fig
        assert fig.axes.ylabel == "Revenue"

    def test_set_legend_returns_figure(self):
        fig = Figure(auto_display=False)
        result = fig.set_legend("bottom-left")
        assert result is fig
        assert fig.legend_pos == "bottom-left"

    def test_set_legend_false(self):
        fig = Figure(auto_display=False)
        fig.set_legend(False)
        assert fig.legend_pos is None

    def test_save_returns_figure(self, tmp_path):
        fig = Figure(auto_display=False)
        fig.add(LineSeries([1, 2], [3, 4]))
        result = fig.save(str(tmp_path / "out.svg"))
        assert result is fig

    def test_full_fluent_chain(self):
        """All mutating methods chain correctly."""
        fig = (
            Figure(auto_display=False)
            .set_title("Chain")
            .set_theme("dark")
            .set_size(800, 500)
            .set_xlabel("X")
            .set_ylabel("Y")
            .set_legend("bottom-right")
            .add(LineSeries([1, 2, 3], [4, 5, 6], label="s"))
            .annotate("hi", x=3, y=6)
        )
        assert isinstance(fig, Figure)
        svg = fig.render_svg()
        assert "Chain" in svg
        assert "#1e1e1e" in svg   # dark background


# ============================================================
# Feature 3 — CLI Tool
# ============================================================

class TestCLI:

    def _run(self, args: list[str]) -> tuple[int, str, str]:
        """Run CLI in-process, capture stderr, return (exit_code, stdout, stderr)."""
        from glyphx.cli import main
        import io
        old_err = sys.stderr
        sys.stderr = io.StringIO()
        try:
            code = main(args)
            err  = sys.stderr.getvalue()
        except SystemExit as e:
            code = e.code
            err  = sys.stderr.getvalue()
        finally:
            sys.stderr = old_err
        return code, "", err

    def _csv(self, tmp_path: Path) -> str:
        """Write a small CSV fixture and return its path."""
        p = tmp_path / "data.csv"
        p.write_text("month,revenue,region\nJan,100,N\nFeb,200,S\nMar,150,N\n")
        return str(p)

    def test_version_subcommand(self):
        code, _, _ = self._run(["version"])
        assert code == 0

    def test_plot_line_to_svg(self, tmp_path):
        csv = self._csv(tmp_path)
        out = str(tmp_path / "chart.svg")
        code, _, _ = self._run(["plot", csv, "--x", "month", "--y", "revenue",
                                 "--kind", "line", "-o", out])
        assert code == 0
        assert os.path.exists(out)
        assert "<svg" in open(out).read()

    def test_plot_bar_to_html(self, tmp_path):
        csv = self._csv(tmp_path)
        out = str(tmp_path / "chart.html")
        code, _, _ = self._run(["plot", csv, "--x", "month", "--y", "revenue",
                                 "--kind", "bar", "-o", out])
        assert code == 0
        assert os.path.exists(out)
        assert "<svg" in open(out).read()

    def test_plot_hist(self, tmp_path):
        csv = self._csv(tmp_path)
        out = str(tmp_path / "hist.svg")
        code, _, _ = self._run(["plot", csv, "--y", "revenue",
                                 "--kind", "hist", "--bins", "5", "-o", out])
        assert code == 0
        assert os.path.exists(out)

    def test_plot_pie(self, tmp_path):
        csv = self._csv(tmp_path)
        out = str(tmp_path / "pie.svg")
        code, _, _ = self._run(["plot", csv, "--x", "region", "--y", "revenue",
                                 "--kind", "pie", "-o", out])
        assert code == 0

    def test_plot_scatter(self, tmp_path):
        csv = self._csv(tmp_path)
        out = str(tmp_path / "scatter.svg")
        code, _, _ = self._run(["plot", csv, "--x", "revenue", "--y", "revenue",
                                 "--kind", "scatter", "-o", out])
        assert code == 0

    def test_plot_with_theme(self, tmp_path):
        csv = self._csv(tmp_path)
        out = str(tmp_path / "dark.html")
        code, _, _ = self._run(["plot", csv, "--y", "revenue",
                                 "--theme", "dark", "-o", out])
        assert code == 0
        assert "#1e1e1e" in open(out).read()

    def test_plot_with_title(self, tmp_path):
        csv = self._csv(tmp_path)
        out = str(tmp_path / "titled.svg")
        code, _, _ = self._run(["plot", csv, "--y", "revenue",
                                 "--title", "My CLI Chart", "-o", out])
        assert code == 0
        assert "My CLI Chart" in open(out).read()

    def test_plot_groupby(self, tmp_path):
        csv = self._csv(tmp_path)
        out = str(tmp_path / "grouped.svg")
        code, _, _ = self._run(["plot", csv, "--y", "revenue",
                                 "--groupby", "region", "--kind", "bar", "-o", out])
        assert code == 0

    def test_plot_missing_file_returns_error(self, tmp_path):
        code, _, _ = self._run(["plot", "/nonexistent/data.csv", "-o", "/tmp/x.svg"])
        assert code != 0

    def test_suggest_subcommand(self, tmp_path):
        csv = self._csv(tmp_path)
        from glyphx.cli import main
        import io
        old_out = sys.stdout
        sys.stdout = io.StringIO()
        try:
            code = main(["suggest", csv])
            out  = sys.stdout.getvalue()
        finally:
            sys.stdout = old_out
        assert code == 0
        assert "bar" in out or "line" in out or "scatter" in out

    def test_suggest_shows_column_names(self, tmp_path):
        csv = self._csv(tmp_path)
        from glyphx.cli import main
        import io
        old_out = sys.stdout
        sys.stdout = io.StringIO()
        try:
            main(["suggest", csv])
            out = sys.stdout.getvalue()
        finally:
            sys.stdout = old_out
        assert "revenue" in out or "month" in out


# ============================================================
# Feature 4 — PPTX Export
# ============================================================

class TestPPTXExport:

    def test_pptx_raises_without_pptx_package(self, tmp_path, basic_fig):
        """If python-pptx / cairosvg not installed, get a clear RuntimeError."""
        out = str(tmp_path / "chart.pptx")
        # If the optional deps ARE installed this test is skipped;
        # if not, confirm the error message is helpful.
        # cairosvg raises OSError (not ImportError) when the system libcairo
        # C library is absent (e.g. on bare macOS CI runners).
        try:
            import pptx   # noqa: F401
            import cairosvg  # noqa: F401
            pytest.skip("python-pptx and cairosvg are installed — skipping error path test")
        except (ImportError, OSError):
            with pytest.raises(RuntimeError, match="pptx|cairosvg|pip"):
                basic_fig.save(out)

    def test_pptx_creates_file_when_deps_available(self, tmp_path, basic_fig):
        """If optional deps are present, a valid .pptx file is produced."""
        try:
            import pptx      # noqa: F401
            import cairosvg  # noqa: F401
        except (ImportError, OSError):
            pytest.skip("python-pptx or cairosvg not installed / libcairo absent")

        out = str(tmp_path / "chart.pptx")
        basic_fig.save(out)

        assert os.path.exists(out)
        assert os.path.getsize(out) > 1000   # not empty

        # Verify it's a valid PPTX (ZIP with expected structure)
        import zipfile
        assert zipfile.is_zipfile(out)
        with zipfile.ZipFile(out) as zf:
            names = zf.namelist()
            assert any("ppt/slides" in n for n in names)

    def test_pptx_with_title(self, tmp_path):
        """Title should appear as a text box on the slide."""
        try:
            import pptx
            import cairosvg
        except (ImportError, OSError):
            pytest.skip("python-pptx or cairosvg not installed / libcairo absent")

        fig = Figure(title="PPTX Title Test", auto_display=False)
        fig.add(LineSeries([1, 2, 3], [4, 5, 6]))
        out = str(tmp_path / "titled.pptx")
        fig.save(out)

        prs = pptx.Presentation(out)
        slide = prs.slides[0]
        texts = [
            shape.text_frame.text
            for shape in slide.shapes
            if shape.has_text_frame
        ]
        assert any("PPTX Title Test" in t for t in texts)

    def test_save_returns_figure_for_pptx(self, tmp_path):
        """fig.save() must return self even for .pptx so chaining works."""
        try:
            import pptx      # noqa: F401
            import cairosvg  # noqa: F401
        except (ImportError, OSError):
            pytest.skip("python-pptx or cairosvg not installed / libcairo absent")

        fig = Figure(auto_display=False)
        fig.add(LineSeries([1, 2], [3, 4]))
        out    = str(tmp_path / "chain.pptx")
        result = fig.save(out)
        assert result is fig


# ============================================================
# Feature 5 — Accessibility
# ============================================================

class TestAccessibility:

    def _svg(self, fig=None):
        if fig is None:
            fig = Figure(title="Access Test", auto_display=False)
            fig.add(LineSeries([1, 2, 3], [4, 5, 6], label="Series A"))
        return fig.render_svg()

    # ── ARIA attributes ─────────────────────────────────────────────────

    def test_svg_has_role_img(self):
        svg = self._svg()
        assert 'role="img"' in svg

    def test_svg_has_aria_labelledby(self):
        svg = self._svg()
        assert "aria-labelledby" in svg

    def test_svg_has_title_element(self):
        svg = self._svg()
        assert "<title" in svg and "</title>" in svg

    def test_svg_has_desc_element(self):
        svg = self._svg()
        assert "<desc" in svg and "</desc>" in svg

    def test_title_content_in_svg(self):
        fig = Figure(title="Revenue Dashboard", auto_display=False)
        fig.add(LineSeries([1, 2], [3, 4]))
        svg = fig.render_svg()
        assert "Revenue Dashboard" in svg

    def test_default_title_when_none(self):
        fig = Figure(auto_display=False)
        fig.add(LineSeries([1, 2], [3, 4]))
        svg = fig.render_svg()
        assert "<title" in svg
        assert "GlyphX Chart" in svg

    def test_title_and_desc_ids_match_aria_labelledby(self):
        svg = self._svg()
        # Find aria-labelledby value
        match = re.search(r'aria-labelledby="([^"]+)"', svg)
        assert match, "aria-labelledby not found"
        ids = match.group(1).split()
        for id_ in ids:
            assert f'id="{id_}"' in svg, f"id '{id_}' referenced but not found in SVG"

    def test_focusable_false_on_svg(self):
        """SVG root must not grab keyboard focus itself — points do."""
        svg = self._svg()
        assert 'focusable="false"' in svg

    # ── tabindex on interactive points ──────────────────────────────────

    def test_glyphx_points_have_tabindex(self):
        svg = self._svg()
        # Every .glyphx-point must be keyboard focusable
        points = re.findall(r'class="[^"]*glyphx-point[^"]*"[^>]*>', svg)
        assert len(points) > 0, "No .glyphx-point elements found"
        for p in points:
            assert 'tabindex="0"' in p, f"tabindex missing on: {p[:80]}"

    def test_glyphx_points_have_role_graphics_symbol(self):
        svg = self._svg()
        points = re.findall(r'class="[^"]*glyphx-point[^"]*"[^>]*>', svg)
        for p in points:
            assert 'role="graphics-symbol"' in p

    def test_bar_points_have_tabindex(self):
        fig = Figure(auto_display=False)
        fig.add(BarSeries(["A", "B", "C"], [1, 2, 3]))
        svg = fig.render_svg()
        assert 'tabindex="0"' in svg

    def test_scatter_points_have_tabindex(self):
        fig = Figure(auto_display=False)
        fig.add(ScatterSeries([1, 2, 3], [4, 5, 6]))
        svg = fig.render_svg()
        assert 'tabindex="0"' in svg

    # ── to_alt_text() ───────────────────────────────────────────────────

    def test_to_alt_text_returns_string(self, basic_fig):
        alt = basic_fig.to_alt_text()
        assert isinstance(alt, str)
        assert len(alt) > 10

    def test_to_alt_text_mentions_chart_kind(self):
        fig = Figure(auto_display=False)
        fig.add(LineSeries([1, 2, 3], [4, 5, 6]))
        alt = fig.to_alt_text()
        assert "line" in alt.lower()

    def test_to_alt_text_includes_title(self):
        fig = Figure(title="Q3 Revenue", auto_display=False)
        fig.add(LineSeries([1, 2], [3, 4]))
        alt = fig.to_alt_text()
        assert "Q3 Revenue" in alt

    def test_to_alt_text_includes_axis_labels(self):
        fig = Figure(auto_display=False)
        fig.add(LineSeries([1, 2], [3, 4]))
        fig.axes.xlabel = "Month"
        fig.axes.ylabel = "USD"
        alt = fig.to_alt_text()
        assert "Month" in alt
        assert "USD" in alt

    def test_to_alt_text_includes_data_range(self):
        fig = Figure(auto_display=False)
        fig.add(LineSeries([1, 2, 3], [10, 50, 30], label="Rev"))
        alt = fig.to_alt_text()
        # Should mention min and max values
        assert "10" in alt and "50" in alt

    def test_to_alt_text_pie_mentions_slices(self):
        from glyphx.series import PieSeries
        fig = Figure(auto_display=False)
        fig.add(PieSeries([30, 40, 30], labels=["A", "B", "C"]))
        alt = fig.to_alt_text()
        assert "slice" in alt.lower() or "3" in alt

    # ── Accessibility JS ────────────────────────────────────────────────

    def test_accessibility_js_exists(self):
        js_path = Path(glyphx.__file__).parent / "assets" / "accessibility.js"
        assert js_path.exists()

    def test_accessibility_js_handles_keyboard_events(self):
        js = (Path(glyphx.__file__).parent / "assets" / "accessibility.js").read_text()
        assert "ArrowRight" in js or "ArrowDown" in js
        assert "Enter" in js
        assert "Escape" in js

    def test_accessibility_js_sets_aria_label(self):
        js = (Path(glyphx.__file__).parent / "assets" / "accessibility.js").read_text()
        assert "aria-label" in js

    def test_accessibility_js_inlined_in_share_html(self, basic_fig):
        html = basic_fig.share()
        assert "ArrowRight" in html or "ArrowDown" in html

    def test_desc_contains_alt_text(self):
        fig = Figure(title="Test", auto_display=False)
        fig.add(LineSeries([1, 2, 3], [10, 50, 30], label="Series"))
        svg = fig.render_svg()
        # Extract <desc> content
        match = re.search(r"<desc[^>]*>([^<]+)</desc>", svg)
        assert match, "<desc> element not found"
        desc_text = match.group(1)
        assert len(desc_text) > 5   # not empty