"""
GlyphX Figure: top-level chart canvas with rendering, display, and export.
"""
from __future__ import annotations

import re
import webbrowser
from tempfile import NamedTemporaryFile
from typing import Any

from .layout import Axes
from .utils import (
    wrap_svg_with_template,
    write_svg_file,
    wrap_svg_canvas,
    draw_legend,
    svg_escape,
)


class Figure:
    """
    Central class for creating and rendering GlyphX visualizations.

    Every mutating method returns ``self`` so calls can be chained::

        fig = (
            Figure(width=900)
            .set_theme("dark")
            .set_title("Monthly Revenue")
            .add(LineSeries(months, revenue, label="Revenue"))
            .annotate("Peak", x="Oct", y=5400)
            .share("report.html")
        )

    Args:
        width:        Canvas width in pixels.
        height:       Canvas height in pixels.
        padding:      Inner margin between canvas edge and plot area.
        title:        Optional title rendered above all plots.
        theme:        Theme name (str) or custom theme dict.
        rows:         Number of subplot rows.
        cols:         Number of subplot columns.
        auto_display: Auto-render when ``plot()`` or ``__repr__`` is called.
        legend:       Legend position string, or ``False`` to suppress.
        xscale:       ``"linear"`` or ``"log"``.
        yscale:       ``"linear"`` or ``"log"``.
    """

    def __init__(
        self,
        width: int = 640,
        height: int = 480,
        padding: int = 50,
        title: str | None = None,
        theme: str | dict[str, Any] | None = None,
        rows: int = 1,
        cols: int = 1,
        auto_display: bool = True,
        legend: str | bool | None = "outside-right",
        xscale: str = "linear",
        yscale: str = "linear",
    ) -> None:
        self.width        = width
        self.height       = height
        self.padding      = padding
        self.title        = title
        self.rows         = rows
        self.cols         = cols
        self.auto_display = auto_display
        self.xscale       = xscale
        self.yscale       = yscale

        from .themes import themes
        self._theme_name: str = (
            theme if isinstance(theme, str) and theme in themes
            else ("custom" if isinstance(theme, dict) else "default")
        )
        self.theme: dict[str, Any] = (
            themes.get(theme, themes["default"])
            if isinstance(theme, str)
            else (theme or themes["default"])
        )

        self.legend_pos: str | None = (
            None if legend in (False, None) else str(legend)
        )

        self.grid: list[list[Axes | None]] = [
            [None] * self.cols for _ in range(self.rows)
        ]
        # When the legend sits to the right of the chart area, shrink the
        # axes so the chart data never overlaps the legend.  The legend
        # occupies the right margin within the same total canvas width.
        from .utils import LEGEND_GUTTER
        # Reserve right gutter for legend regardless of position string.
        # This prevents legends from ever overlapping chart data.
        _axes_width = self.width - LEGEND_GUTTER if legend not in (False, None) else self.width
        self.axes = Axes(
            width=_axes_width,
            height=self.height,
            padding=self.padding,
            theme=self.theme,
            xscale=xscale,
            yscale=yscale,
        )
        self.series: list[tuple[Any, bool]] = []
        self._annotations:   list[dict[str, Any]] = []
        self._canvas_texts:  list[dict[str, Any]] = []
        self._supxlabel:     dict | None = None
        self._supylabel:     dict | None = None

    # -- Fluent setters ---------------------------------------------------

    def set_title(self, title: str) -> Figure:
        """Set the figure title and return ``self`` for chaining."""
        self.title = title
        return self

    def set_theme(self, theme: str | dict[str, Any]) -> Figure:
        """Apply a named theme or a custom dict and return ``self``."""
        from .themes import themes
        self._theme_name = (
            theme if isinstance(theme, str) and theme in themes
            else ("custom" if isinstance(theme, dict) else "default")
        )
        self.theme = (
            themes.get(theme, themes["default"])
            if isinstance(theme, str)
            else theme
        )
        self.axes.theme = self.theme
        return self

    def set_size(self, width: int, height: int) -> Figure:
        """Resize the canvas and return ``self``."""
        self.width  = width
        self.height = height
        self.axes.width  = width
        self.axes.height = height
        return self

    def set_xlabel(self, label: str) -> Figure:
        """Set the X-axis label and return ``self``."""
        self.axes.xlabel = label
        return self

    def set_ylabel(self, label: str) -> Figure:
        """Set the Y-axis label and return ``self``."""
        self.axes.ylabel = label
        return self

    def set_legend(self, position: str | bool | None) -> Figure:
        """Set legend position (or ``False`` to hide) and return ``self``."""
        self.legend_pos = None if position in (False, None) else str(position)
        return self

    # -- Subplot grid -----------------------------------------------------

    def add_axes(self, row: int = 0, col: int = 0) -> Axes:
        """Create or retrieve the Axes at a grid position."""
        if self.grid[row][col] is None:
            self.grid[row][col] = Axes(
                width=self.width  // self.cols,
                height=self.height // self.rows,
                padding=self.padding,
                theme=self.theme,
                xscale=self.xscale,
                yscale=self.yscale,
            )
        return self.grid[row][col]  # type: ignore[return-value]

    # -- Series management -------------------------------------------------

    def add(self, series: Any, use_y2: bool = False) -> Figure:
        """
        Add a series to the figure.

        Returns ``self`` so calls can be chained::

            fig.add(LineSeries(...)).add(BarSeries(...)).show()
        """
        self.series.append((series, use_y2))
        if hasattr(series, "x") and hasattr(series, "y"):
            self.axes.add_series(series, use_y2)
        return self


    # -- Shorthand series methods ------------------------------------------
    # These mirror the long-form Figure().add(XxxSeries(...)) API so users
    # can write fig.line(x, y) instead of importing and constructing manually.
    # Every method returns self for chaining and accepts the same kwargs as
    # the underlying series class.

    def line(self, x, y, color=None, label=None, linestyle="solid",
             width=2, yerr=None, xerr=None, use_y2=False, **kwargs) -> "Figure":
        """Add a :class:`~glyphx.series.LineSeries`.  Returns ``self``."""
        from .series import LineSeries
        return self.add(LineSeries(x, y, color=color, label=label,
                                   linestyle=linestyle, width=width,
                                   yerr=yerr, xerr=xerr, **kwargs), use_y2)

    def bar(self, x, y, color=None, label=None, bar_width=0.8,
            yerr=None, use_y2=False, **kwargs) -> "Figure":
        """Add a :class:`~glyphx.series.BarSeries`.  Returns ``self``."""
        from .series import BarSeries
        return self.add(BarSeries(x, y, color=color, label=label,
                                  bar_width=bar_width, yerr=yerr, **kwargs), use_y2)

    def scatter(self, x, y, color=None, label=None, size=5,
                c=None, cmap="viridis", use_y2=False, **kwargs) -> "Figure":
        """Add a :class:`~glyphx.series.ScatterSeries`.  Returns ``self``."""
        from .series import ScatterSeries
        return self.add(ScatterSeries(x, y, color=color, label=label,
                                      size=size, c=c, cmap=cmap, **kwargs), use_y2)

    def hist(self, data, bins=20, color=None, label=None, **kwargs) -> "Figure":
        """Add a :class:`~glyphx.series.HistogramSeries`.  Returns ``self``."""
        from .series import HistogramSeries
        return self.add(HistogramSeries(data, bins=bins,
                                        color=color, label=label, **kwargs))

    def box(self, data, categories=None, color=None, box_width=20,
            **kwargs) -> "Figure":
        """Add a :class:`~glyphx.series.BoxPlotSeries`.  Returns ``self``."""
        from .series import BoxPlotSeries
        return self.add(BoxPlotSeries(data, categories=categories,
                                      color=color or "#1f77b4",
                                      box_width=box_width, **kwargs))

    def heatmap(self, matrix, row_labels=None, col_labels=None,
                show_values=False, cmap=None, **kwargs) -> "Figure":
        """Add a :class:`~glyphx.series.HeatmapSeries`.  Returns ``self``."""
        from .series import HeatmapSeries
        return self.add(HeatmapSeries(matrix, row_labels=row_labels,
                                      col_labels=col_labels,
                                      show_values=show_values,
                                      cmap=cmap, **kwargs))

    def pie(self, values, labels=None, colors=None, **kwargs) -> "Figure":
        """Add a :class:`~glyphx.series.PieSeries`.  Returns ``self``."""
        from .series import PieSeries
        return self.add(PieSeries(values=values, labels=labels,
                                  colors=colors, **kwargs))

    def donut(self, values, labels=None, colors=None, **kwargs) -> "Figure":
        """Add a :class:`~glyphx.series.DonutSeries`.  Returns ``self``."""
        from .series import DonutSeries
        return self.add(DonutSeries(values=values, labels=labels,
                                    colors=colors, **kwargs))

    def area(self, x, y1, y2=0, color=None, alpha=0.25,
             line_width=1, label=None, **kwargs) -> "Figure":
        """Add a :class:`~glyphx.fill_between.FillBetweenSeries`.  Returns ``self``."""
        from .fill_between import FillBetweenSeries
        return self.add(FillBetweenSeries(x, y1, y2, color=color or "#1f77b4",
                                          alpha=alpha, line_width=line_width,
                                          label=label, **kwargs))

    def kde(self, data, filled=False, alpha=0.20, color=None,
            width=2, label=None, **kwargs) -> "Figure":
        """Add a :class:`~glyphx.kde.KDESeries`.  Returns ``self``."""
        from .kde import KDESeries
        return self.add(KDESeries(data, filled=filled, alpha=alpha,
                                  color=color or "#1f77b4", width=width,
                                  label=label, **kwargs))

    def ecdf(self, data, color=None, label=None,
             complementary=False, **kwargs) -> "Figure":
        """Add a :class:`~glyphx.ecdf.ECDFSeries`.  Returns ``self``."""
        from .ecdf import ECDFSeries
        return self.add(ECDFSeries(data, color=color, label=label,
                                   complementary=complementary, **kwargs))

    def raincloud(self, data, categories=None, seed=42,
                  violin_width=40, **kwargs) -> "Figure":
        """Add a :class:`~glyphx.raincloud.RaincloudSeries`.  Returns ``self``."""
        from .raincloud import RaincloudSeries
        series = RaincloudSeries(data, categories=categories,
                                 seed=seed, violin_width=violin_width, **kwargs)
        series._x_categories = series.categories
        return self.add(series)

    def candlestick(self, dates, open, high, low, close,
                    label=None, **kwargs) -> "Figure":
        """Add a :class:`~glyphx.candlestick.CandlestickSeries`.  Returns ``self``."""
        from .candlestick import CandlestickSeries
        return self.add(CandlestickSeries(dates, open, high, low, close,
                                          label=label, **kwargs))

    def waterfall(self, labels, values, show_values=True, **kwargs) -> "Figure":
        """Add a :class:`~glyphx.waterfall.WaterfallSeries`.  Returns ``self``."""
        from .waterfall import WaterfallSeries
        return self.add(WaterfallSeries(labels=labels, values=values,
                                        show_values=show_values, **kwargs))

    def treemap(self, labels, values, cmap="viridis",
                show_values=True, **kwargs) -> "Figure":
        """Add a :class:`~glyphx.treemap.TreemapSeries`.  Returns ``self``."""
        from .treemap import TreemapSeries
        return self.add(TreemapSeries(labels=labels, values=values,
                                      cmap=cmap, show_values=show_values,
                                      **kwargs))

    def bubble(self, x, y, size, color=None, c=None, cmap="viridis",
               alpha=0.65, min_radius=4, max_radius=40,
               labels=None, label=None, use_y2=False, **kwargs) -> "Figure":
        """Add a :class:`~glyphx.bubble.BubbleSeries`.  Returns ``self``."""
        from .bubble import BubbleSeries
        return self.add(BubbleSeries(x, y, size, color=color, c=c, cmap=cmap,
                                     alpha=alpha, min_radius=min_radius,
                                     max_radius=max_radius, labels=labels,
                                     label=label, **kwargs), use_y2)

    def sunburst(self, labels, parents, values, cmap="viridis",
                 **kwargs) -> "Figure":
        """Add a :class:`~glyphx.sunburst.SunburstSeries`.  Returns ``self``."""
        from .sunburst import SunburstSeries
        return self.add(SunburstSeries(labels=labels, parents=parents,
                                       values=values, cmap=cmap, **kwargs))

    def parallel_coords(self, data, axes, hue=None, cmap="viridis",
                        alpha=0.45, **kwargs) -> "Figure":
        """Add a :class:`~glyphx.parallel_coords.ParallelCoordinatesSeries`.
        Returns ``self``."""
        from .parallel_coords import ParallelCoordinatesSeries
        return self.add(ParallelCoordinatesSeries(data=data, axes=axes,
                                                  hue=hue, cmap=cmap,
                                                  alpha=alpha, **kwargs))

    def diverging_bar(self, categories, values, pos_color="#2563eb",
                      neg_color="#dc2626", show_values=True,
                      **kwargs) -> "Figure":
        """Add a :class:`~glyphx.diverging_bar.DivergingBarSeries`.
        Returns ``self``."""
        from .diverging_bar import DivergingBarSeries
        return self.add(DivergingBarSeries(categories=categories, values=values,
                                           pos_color=pos_color, neg_color=neg_color,
                                           show_values=show_values, **kwargs))

    def stream(self, max_points=100, color=None, label=None,
               **kwargs) -> "Figure":
        """Add a :class:`~glyphx.streaming.StreamingSeries` and return it.

        Unlike other shorthand methods, this returns the *series* (not self)
        so callers can push data to it::

            stream = fig.stream(max_points=100, label="Sensor")
            stream.push(42.0)
        """
        from .streaming import StreamingSeries
        s = StreamingSeries(max_points=max_points,
                            color=color or "#7c3aed",
                            label=label, **kwargs)
        self.add(s)
        return s

    def axhspan(self, ymin: float, ymax: float,
                color: str = "#facc15", alpha: float = 0.20,
                label: str | None = None) -> "Figure":
        """
        Add a horizontal shaded band across the full chart width.

        Equivalent to Matplotlib's ``ax.axhspan()``.  Y values are in
        data space; the band is clipped to the plot area.

        Args:
            ymin:  Lower Y bound (data units).
            ymax:  Upper Y bound (data units).
            color: Fill color.
            alpha: Fill opacity 0-1.
            label: Optional legend label.

        Returns:
            ``self`` for chaining.

        Example::

            fig.axhspan(90, 110, color="#22c55e", alpha=0.15, label="Target range")
        """
        self.axes.axhspan(ymin, ymax, color=color, alpha=alpha, label=label)
        return self

    def axvspan(self, xmin, xmax,
                color: str = "#a855f7", alpha: float = 0.20,
                label: str | None = None) -> "Figure":
        """
        Add a vertical shaded band across the full chart height.

        Equivalent to Matplotlib's ``ax.axvspan()``.  X values may be
        numeric data values or category label strings.

        Args:
            xmin:  Left X bound (data units or category label).
            xmax:  Right X bound.
            color: Fill color.
            alpha: Fill opacity 0-1.
            label: Optional legend label.

        Returns:
            ``self`` for chaining.

        Example::

            fig.axvspan("Jul", "Sep", color="#f59e0b", alpha=0.15, label="Q3")
        """
        self.axes.axvspan(xmin, xmax, color=color, alpha=alpha, label=label)
        return self

    def set_xticks(self, ticks: list,
                   labels: list | None = None) -> "Figure":
        """
        Set explicit X-axis tick positions (and optionally their labels).

        Equivalent to Matplotlib's ``ax.set_xticks()`` + ``ax.set_xticklabels()``.

        Returns:
            ``self`` for chaining.

        Example::

            fig.set_xticks([0, 25, 50, 75, 100])
            fig.set_xticks([1, 6, 12], labels=["Jan", "Jun", "Dec"])
        """
        self.axes.set_xticks(ticks, labels)
        return self

    def set_yticks(self, ticks: list,
                   labels: list | None = None) -> "Figure":
        """
        Set explicit Y-axis tick positions (and optionally their labels).

        Returns:
            ``self`` for chaining.

        Example::

            fig.set_yticks([0, 500_000, 1_000_000], labels=["$0", "$500k", "$1M"])
        """
        self.axes.set_yticks(ticks, labels)
        return self

    def set_tick_format(self, formatter) -> "Figure":
        """
        Apply a callable formatter to all numeric tick labels on both axes.

        Equivalent to Matplotlib's ``FuncFormatter``.

        Args:
            formatter: Callable ``(value: float) -> str``.

        Returns:
            ``self`` for chaining.

        Example::

            fig.set_tick_format(lambda v: f"${v:,.0f}")
            fig.set_tick_format(lambda v: f"{v:.1%}")
        """
        self.axes.set_tick_format(formatter)
        return self

    def set_minor_ticks(self, n: int = 4) -> "Figure":
        """
        Draw minor tick subdivisions between each pair of major ticks.

        Args:
            n: Number of minor ticks between major ticks.

        Returns:
            ``self`` for chaining.

        Example::

            fig.set_minor_ticks(4)
        """
        self.axes.set_minor_ticks(n)
        return self

    def set_spine_visible(self, left: bool = True, right: bool = True,
                          top: bool = False, bottom: bool = True) -> "Figure":
        """
        Control axis spine (border line) visibility.

        Equivalent to Matplotlib's ``ax.spines[...].set_visible()``.

        Args:
            left:   Left (Y) spine.
            right:  Right (Y2) spine.
            top:    Top spine.
            bottom: Bottom (X) spine.

        Returns:
            ``self`` for chaining.

        Example::

            fig.set_spine_visible(top=False, right=False)   # clean minimal look
        """
        self.axes.set_spine_visible(left=left, right=right,
                                     top=top, bottom=bottom)
        return self

    def text(self, x: float, y: float, s: str,
             color: str = "#333", font_size: int = 12,
             anchor: str = "middle",
             transform: str = "canvas") -> "Figure":
        """
        Add free-form text anywhere on the canvas.

        Equivalent to Matplotlib's ``fig.text()`` / ``ax.text()``.

        Args:
            x:         X position.  When ``transform='canvas'`` this is
                       a fraction of canvas width (0-1).  When
                       ``transform='data'`` it is a data-space value.
            y:         Y position.  Canvas fraction (0 = top, 1 = bottom)
                       or data-space value depending on ``transform``.
            s:         Text string.  Use ``$...$`` for math (rendered
                       via MathJax when the chart is opened in a browser).
            color:     Font color.
            font_size: Font size in points.
            anchor:    SVG text-anchor: ``"start"``, ``"middle"``, ``"end"``.
            transform: ``"canvas"`` (0-1 fractions) or ``"data"`` (data coords).

        Returns:
            ``self`` for chaining.

        Example::

            fig.text(0.5, 0.02, "Source: World Bank", font_size=9, color="#888")
            fig.text(0.5, 0.5, "$R^2 = 0.95$", font_size=14)
        """
        self._canvas_texts.append(dict(
            x=x, y=y, s=s, color=color,
            font_size=font_size, anchor=anchor, transform=transform,
        ))
        return self

    def supxlabel(self, label: str, font_size: int = 13,
                  color: str | None = None) -> "Figure":
        """
        Set a shared X-axis label below all subplots.

        Equivalent to Matplotlib's ``fig.supxlabel()``.  Useful when all
        subplots in a grid share the same X-axis meaning.

        Args:
            label:     Label text.
            font_size: Font size.
            color:     Text color (defaults to theme text color).

        Returns:
            ``self`` for chaining.
        """
        self._supxlabel = dict(label=label, font_size=font_size, color=color)
        return self

    def supylabel(self, label: str, font_size: int = 13,
                  color: str | None = None) -> "Figure":
        """
        Set a shared Y-axis label beside all subplots.

        Equivalent to Matplotlib's ``fig.supylabel()``.

        Args:
            label:     Label text.
            font_size: Font size.
            color:     Text color (defaults to theme text color).

        Returns:
            ``self`` for chaining.
        """
        self._supylabel = dict(label=label, font_size=font_size, color=color)
        return self

    def to_vega_lite(self, path: str | None = None, **kwargs) -> dict:
        """
        Convert this figure to a Vega-Lite v5 specification.

        The spec can be saved as a ``.vl.json`` file, embedded in
        Observable notebooks, or rendered with ``altair``.
        No Python plotting library currently produces Vega-Lite output.

        Args:
            path: If given, write the spec to this JSON file.

        Returns:
            Vega-Lite v5 specification dict.

        Example::

            spec = fig.to_vega_lite()
            fig.to_vega_lite("chart.vl.json")  # saves to file
        """
        from .vega_lite import to_vega_lite as _to_vl, save_vega_lite
        spec = _to_vl(self, **kwargs)
        if path is not None:
            save_vega_lite(self, path, **kwargs)
        return spec

    def regplot(self, data=None, x=None, y=None,
                x_vals=None, y_vals=None, **kwargs) -> 'Figure':
        """Add a regression plot series to this figure. Returns ``self``."""
        from .regplot import regplot as _regplot
        reg_fig = _regplot(data, x=x, y=y,
                           x_vals=x_vals, y_vals=y_vals,
                           auto_display=False, **kwargs)
        for s, u in reg_fig.series:
            self.add(s, y2=u)
        return self

    def choropleth(self, geojson, data, key='name', cmap='viridis',
                   **kwargs) -> 'Figure':
        """Add a :class:`~glyphx.choropleth.ChoroplethSeries`. Returns ``self``."""
        from .choropleth import ChoroplethSeries
        return self.add(ChoroplethSeries(geojson, data, key=key,
                                          cmap=cmap, **kwargs))

    def gantt(self, tasks: list, group_colors: dict | None = None,
              cmap: str = 'viridis', bar_height: int = 20,
              show_today: bool = True, **kwargs) -> 'Figure':
        """Add a :class:`~glyphx.gantt.GanttSeries`.  Returns ``self``."""
        from .gantt import GanttSeries
        return self.add(GanttSeries(tasks=tasks, group_colors=group_colors,
                                     cmap=cmap, bar_height=bar_height,
                                     show_today=show_today, **kwargs))

    def stacked_bar(self, x: list, series: dict,
                    normalize: bool = False, colors: list | None = None,
                    bar_width: float = 0.75, **kwargs) -> "Figure":
        """Add a :class:`~glyphx.stacked_bar.StackedBarSeries`.  Returns ``self``."""
        from .stacked_bar import StackedBarSeries
        return self.add(StackedBarSeries(x=x, series=series,
                                          normalize=normalize, colors=colors,
                                          bar_width=bar_width, **kwargs))

    def bump(self, x: list, rankings: dict,
             colors: list | None = None, show_labels: bool = True,
             **kwargs) -> "Figure":
        """Add a :class:`~glyphx.bump_chart.BumpChartSeries`.  Returns ``self``."""
        from .bump_chart import BumpChartSeries
        return self.add(BumpChartSeries(x=x, rankings=rankings,
                                         colors=colors, show_labels=show_labels,
                                         **kwargs))

    def sparkline(self, data: list, color: str = "#2563eb",
                  kind: str = "line", fill: bool = True,
                  label: str | None = None, **kwargs) -> "Figure":
        """Add a :class:`~glyphx.sparkline.SparklineSeries`.  Returns ``self``."""
        from .sparkline import SparklineSeries
        return self.add(SparklineSeries(data=data, color=color,
                                         kind=kind, fill=fill,
                                         label=label, **kwargs))

    def vline(self, x, color="#888", width=1,
              linestyle="dashed", label=None) -> "Figure":
        """Draw a vertical reference line at data coordinate ``x``.  Returns ``self``."""
        from .series import LineSeries
        # Compute domain from existing series if not yet finalized
        if self.axes._y_domain:
            ymin, ymax = self.axes._y_domain
        elif self.series:
            all_y = [v for s, _ in self.series for v in (s.y or []) if v is not None]
            ymin, ymax = (min(all_y), max(all_y)) if all_y else (0, 1)
        else:
            ymin, ymax = 0, 1
        return self.add(LineSeries([x, x], [ymin, ymax],
                                   color=color, width=width,
                                   linestyle=linestyle, label=label))

    def hline(self, y, color="#888", width=1,
              linestyle="dashed", label=None) -> "Figure":
        """Draw a horizontal reference line at data coordinate ``y``.  Returns ``self``."""
        from .series import LineSeries
        if self.axes._x_domain:
            xmin, xmax = self.axes._x_domain
        elif self.series:
            all_x = [v for s, _ in self.series for v in (getattr(s, "_numeric_x", None) or s.x or []) if v is not None]
            xmin, xmax = (min(all_x), max(all_x)) if all_x else (0, 1)
        else:
            xmin, xmax = 0, 1
        return self.add(LineSeries([xmin, xmax], [y, y],
                                   color=color, width=width,
                                   linestyle=linestyle, label=label))

    # -- __repr__ ---------------------------------------------------------

    def copy(self) -> "Figure":
        """
        Return a deep copy of this figure.

        Useful for creating small variations without re-building from scratch::

            base = Figure().add(LineSeries(x, y))
            dark = base.copy().set_theme("dark")
            light = base.copy().set_theme("default")

        Returns:
            A new :class:`Figure` with all series, annotations, and settings
            duplicated.
        """
        import copy as _copy
        return _copy.deepcopy(self)

    def __eq__(self, other: object) -> bool:
        """
        Compare two figures by their rendered SVG output.

        Chart IDs (UUIDs) are stripped before comparison so two figures
        produced from identical data are considered equal even though their
        IDs differ.  This enables snapshot / regression testing::

            assert fig_a == fig_b           # visual equality
            assert fig_new != fig_baseline  # regression detected

        Returns:
            ``True`` if both figures render to visually identical SVG.
        """
        if not isinstance(other, Figure):
            return NotImplemented
        import re as _re
        # Pattern strips any XML attribute whose value contains a glyphx UUID
        _UUID_ATTR = _re.compile(r'\s[\w:-]+="[^"]*glyphx-chart-[a-f0-9]{12}[^"]*"')

        def _strip_ids(s: str) -> str:
            return _UUID_ATTR.sub("", s)

        return _strip_ids(self.render_svg()) == _strip_ids(other.render_svg())

    def __hash__(self) -> int:
        """
        Hash based on figure identity (not SVG content).

        Required by Python when ``__eq__`` is defined.  Figures remain
        usable as dict keys and in sets even though equality is SVG-based.
        """
        return id(self)

    def __repr__(self) -> str:
        series_desc = ", ".join(
            f"{s.__class__.__name__}({repr(s.label)})"
            if getattr(s, "label", None)
            else s.__class__.__name__
            for s, _ in self.series
        )
        theme_name = getattr(self, "_theme_name", "default")
        return (
            f"<glyphx.Figure {self.width}x{self.height}"
            + (f" [{series_desc}]" if self.series else " [empty]")
            + f" theme={theme_name!r}>"
        )

    # -- Annotations ------------------------------------------------------

    def annotate(
        self,
        text: str,
        x: float,
        y: float,
        color: str = "#333",
        font_size: int = 12,
        anchor: str = "start",
        arrow: bool = False,
        ax_x: float | None = None,
        ax_y: float | None = None,
    ) -> Figure:
        """
        Add a text annotation in data-space coordinates.

        Returns ``self`` for chaining.

        Args:
            text:      Label text.
            x:         Data-space X coordinate of the annotation point.
            y:         Data-space Y coordinate.
            color:     Text and arrow colour.
            font_size: Label font size in points.
            anchor:    SVG text-anchor (``"start"``, ``"middle"``, ``"end"``).
            arrow:     Draw a small leader line from text to point.
            ax_x:      Arrow tail X pixel offset (default −8).
            ax_y:      Arrow tail Y pixel offset (default −8).
        """
        self._annotations.append(dict(
            text=text, x=x, y=y, color=color,
            font_size=font_size, anchor=anchor,
            arrow=arrow, ax_x=ax_x, ax_y=ax_y,
        ))
        return self

    def _render_annotations(
        self,
        scale_x: Any,
        scale_y: Any,
        font: str,
    ) -> str:
        elements: list[str] = []
        # Build a category->numeric map from all registered series
        cat_map: dict = {}
        for s, _ in self.series:
            if hasattr(s, "_x_categories") and s._x_categories:
                for i, cat in enumerate(s._x_categories):
                    cat_map[str(cat)] = i + 0.5

        for ann in self._annotations:
            ann_x = ann["x"]
            ann_y = ann["y"]
            # Resolve categorical x values to numeric
            if isinstance(ann_x, str) and ann_x in cat_map:
                ann_x = cat_map[ann_x]
            try:
                px  = scale_x(float(ann_x))
            except (TypeError, ValueError):
                continue
            try:
                py  = scale_y(float(ann_y))
            except (TypeError, ValueError):
                continue
            ox  = ann["ax_x"] if ann["ax_x"] is not None else -8
            oy  = ann["ax_y"] if ann["ax_y"] is not None else -8
            if ann["arrow"]:
                elements.append(
                    f'<line x1="{px + ox}" y1="{py + oy}" x2="{px}" y2="{py}" '
                    f'stroke="{ann["color"]}" stroke-width="1.5" '
                    f'marker-end="url(#arrow)"/>'
                )
            elements.append(
                f'<text x="{px + ox}" y="{py + oy - 2}" '
                f'text-anchor="{ann["anchor"]}" font-size="{ann["font_size"]}" '
                f'font-family="{font}" fill="{ann["color"]}">'
                f'{svg_escape(ann["text"])}</text>'
            )
        return "\n".join(elements)

    @staticmethod
    def _arrow_marker_def() -> str:
        return (
            '<defs><marker id="arrow" markerWidth="8" markerHeight="8" '
            'refX="6" refY="3" orient="auto">'
            '<path d="M0,0 L0,6 L8,3 z" fill="#333"/>'
            '</marker></defs>'
        )

    # -- Accessibility -----------------------------------------------------

    def to_alt_text(self) -> str:
        """
        Generate a plain-English description of this figure for screen readers.

        Returns:
            A human-readable string suitable for ``aria-label`` or ``<desc>``.
        """
        from .a11y import generate_alt_text
        return generate_alt_text(self)

    # -- Rendering --------------------------------------------------------

    def render_svg(self, viewbox: bool = False) -> str:
        """
        Render the complete figure and return an SVG string.

        The SVG includes:
        - ``role="img"`` and ``aria-labelledby`` on the root element
        - ``<title>`` and ``<desc>`` ARIA landmark children
        - ``tabindex="0"`` on every interactive data point

        Returns:
            Complete SVG document markup.
        """
        svg_parts: list[str] = []

        if any(a["arrow"] for a in self._annotations):
            svg_parts.append(self._arrow_marker_def())

        svg_parts.append(
            f'<rect width="{self.width}" height="{self.height}" '
            f'fill="{self.theme.get("background", "#ffffff")}"/>'
        )

        if self.title:
            font  = self.theme.get("font", "sans-serif")
            color = self.theme.get("text_color", "#000")
            svg_parts.append(
                f'<text x="{self.width // 2}" y="28" text-anchor="middle" '
                f'font-size="20" font-weight="bold" font-family="{font}" '
                f'fill="{color}">{svg_escape(self.title)}</text>'
            )

        # sup[x|y]label -- shared axis labels across subplot grids
        _font  = self.theme.get("font", "sans-serif")
        _tc    = self.theme.get("text_color", "#000")
        if self._supxlabel:
            sx = self._supxlabel
            c  = sx.get("color") or _tc
            svg_parts.append(
                f'<text x="{self.width // 2}" y="{self.height - 4}" '
                f'text-anchor="middle" font-size="{sx["font_size"]}" '
                f'font-family="{_font}" fill="{c}">{svg_escape(sx["label"])}</text>'
            )
        if self._supylabel:
            sy = self._supylabel
            c  = sy.get("color") or _tc
            cx = 14
            cy = self.height // 2
            svg_parts.append(
                f'<text x="{cx}" y="{cy}" text-anchor="middle" '
                f'font-size="{sy["font_size"]}" font-family="{_font}" fill="{c}" '
                f'transform="rotate(-90, {cx}, {cy})">{svg_escape(sy["label"])}</text>'
            )

        # Free-form canvas text (fig.text())
        for ct in getattr(self, "_canvas_texts", []):
            if ct["transform"] == "canvas":
                cx = ct["x"] * self.width
                cy = ct["y"] * self.height
            else:
                # data coords -- resolve if axes are finalized
                try:
                    cx = self.axes.scale_x(ct["x"]) if self.axes.scale_x else ct["x"] * self.width
                    cy = self.axes.scale_y(ct["y"]) if self.axes.scale_y else ct["y"] * self.height
                except Exception:
                    cx, cy = ct["x"] * self.width, ct["y"] * self.height
            svg_parts.append(
                f'<text x="{cx:.1f}" y="{cy:.1f}" '
                f'text-anchor="{ct["anchor"]}" '
                f'font-size="{ct["font_size"]}" font-family="{_font}" '
                f'fill="{ct["color"]}">{svg_escape(ct["s"])}</text>'
            )

        # -- Subplot grid --------------------------------------------------
        if self.grid and any(any(cell for cell in row) for row in self.grid):
            cell_w = self.width  // self.cols
            cell_h = self.height // self.rows
            for r, row in enumerate(self.grid):
                for c, ax in enumerate(row):
                    if not ax:
                        continue
                    ax.finalize()
                    group = f'<g transform="translate({c * cell_w},{r * cell_h})">'
                    group += ax.render_axes() + ax.render_grid()
                    for s in ax.series:
                        group += s.to_svg(ax)
                    if getattr(ax, "legend_pos", None):
                        group += draw_legend(
                            ax.series,
                            position=ax.legend_pos,
                            font=self.theme.get("font", "sans-serif"),
                            text_color=self.theme.get("text_color", "#000"),
                            fig_width=ax.width,
                            fig_height=ax.height,
                        )
                    group += "</g>"
                    svg_parts.append(group)

        # -- Single-axes ---------------------------------------------------
        elif self.series and any(
            hasattr(s, "x") and hasattr(s, "y") and s.x and s.y
            for s, _ in self.series
        ):
            if not self.axes.series:
                for s, use_y2 in self.series:
                    self.axes.add_series(s, use_y2)

            self.axes.finalize()
            svg_parts.append(self.axes.render_axes())
            svg_parts.append(self.axes.render_grid())

            for series, use_y2 in self.series:
                svg_parts.append(series.to_svg(self.axes, use_y2=use_y2))

            if self._annotations and self.axes.scale_x and self.axes.scale_y:
                svg_parts.append(self._render_annotations(
                    self.axes.scale_x,
                    self.axes.scale_y,
                    self.theme.get("font", "sans-serif"),
                ))

            if self.legend_pos:
                # For outside-right legends, anchor x relative to the
                # axes width (the shrunk plot area) not the full canvas.
                # Always render legend in the right gutter (outside plot area)
                _legend_ref_w = self.axes.width
                svg_parts.append(draw_legend(
                    [s for s, _ in self.series],
                    position="outside-right",
                    font=self.theme.get("font", "sans-serif"),
                    text_color=self.theme.get("text_color", "#000"),
                    fig_width=_legend_ref_w,
                    fig_height=self.height,
                ))

            # -- Statistical annotations -----------------------------------
            for ann in getattr(self, "_stat_annotations", []):
                svg_parts.append(ann.to_svg(self.axes))

        # -- Axis-free (pie, donut, etc.) ----------------------------------
        elif self.series:
            for series, _ in self.series:
                svg_parts.append(series.to_svg(self.axes))

        # Detect math text ($...$) in the rendered SVG content for MathJax
        _svg_content = "\n".join(svg_parts)
        _has_math    = "$" in _svg_content

        raw_svg = wrap_svg_canvas(
            _svg_content,
            width=self.width,
            height=self.height,
            has_math=_has_math,
        )

        # -- Accessibility injection ---------------------------------------
        chart_id = re.search(r'id="(glyphx-chart-[^"]+)"', raw_svg)
        cid      = chart_id.group(1) if chart_id else "glyphx-chart-0"

        from .a11y import inject_aria
        return inject_aria(
            svg=raw_svg,
            title=self.title or "GlyphX Chart",
            desc=self.to_alt_text(),
            chart_id=cid,
        )

    # -- Display / export --------------------------------------------------

    def _display(self, svg_string: str) -> None:
        try:
            from IPython import get_ipython
            ip = get_ipython()
            if ip is not None and "IPKernelApp" in ip.config:
                from IPython.display import SVG, display as jup_display
                jup_display(SVG(svg_string))
                return
        except Exception:
            pass
        html = wrap_svg_with_template(svg_string)
        tmp  = NamedTemporaryFile(delete=False, suffix=".html", mode="w", encoding="utf-8")
        tmp.write(html)
        tmp.close()
        webbrowser.open(f"file://{tmp.name}")

    def show(self) -> Figure:
        """Render and display the figure. Returns ``self`` for chaining."""
        self._display(self.render_svg())
        return self

    def render_responsive(self,
                          dark_theme: str | None = None) -> str:
        """
        Render an SVG that responds to the OS dark-mode preference.

        The output uses CSS custom properties (``--glyphx-bg``, etc.).
        When the user switches dark mode, the chart recolours without
        any Python re-render -- a capability none of the three competitors
        have.

        Args:
            dark_theme: GlyphX theme name for dark mode (default ``'dark'``).

        Returns:
            Complete ``<svg>`` string with embedded CSS variable block.

        Example::

            svg = fig.render_responsive(dark_theme='dark')
            Path('chart.svg').write_text(svg)
            # Open in a browser -- switches colours with the OS setting
        """
        from .utils   import wrap_svg_with_css_vars
        from .themes  import themes as _themes

        light  = self.theme
        dark   = _themes.get(dark_theme or 'dark', _themes['dark'])

        # Render inner SVG content (no root element)
        svg_parts: list[str] = []
        if self.title:
            font  = light.get('font', 'sans-serif')
            svg_parts.append(
                f'<text x="{self.width // 2}" y="28" text-anchor="middle" '
                f'font-size="20" font-weight="bold" font-family="{font}" '
                f'fill="var(--glyphx-text)">{svg_escape(self.title)}</text>'
            )
        inner = '\n'.join(svg_parts)
        return wrap_svg_with_css_vars(
            inner + '\n' + self.render_svg(),
            light_theme=light, dark_theme=dark,
            width=self.width, height=self.height,
        )

    def save(self, filename: str = "glyphx_output.svg",
             dpi: int = 96) -> "Figure":
        """
        Save the rendered figure to disk.

        Supported extensions: ``.svg``, ``.html``, ``.png``, ``.jpg``,
        ``.pptx``.  PNG/JPG/PPTX require optional extras::

            pip install "glyphx[export]"    # PNG/JPG
            pip install "glyphx[pptx]"      # PowerPoint

        Args:
            filename: Output path.  Extension determines the format.
            dpi:      Output resolution for raster formats (PNG/JPG).
                      Default 96; use 192 or 300 for high-DPI / print.
                      Has no effect on SVG or HTML output.

        Returns:
            ``self`` for chaining.

        Example::

            fig.save("chart.png", dpi=192)   # crisp on retina displays
            fig.save("chart.png", dpi=300)   # print-quality
        """
        svg = self.render_svg()
        if filename.lower().endswith(".pptx"):
            _save_as_pptx(svg, filename, title=self.title)
        else:
            write_svg_file(svg, filename, dpi=dpi)
        return self


    def tight_layout(self) -> "Figure":
        """
        Auto-adjust padding to prevent label clipping and overlap.

        For single-axis figures, delegates to
        :meth:`~glyphx.layout.Axes.tight_layout`.  For subplot grids,
        applies :meth:`constrained_layout` which aligns the left edges of
        all subplot cells so tick labels never overlap across rows.

        Returns ``self`` for chaining.
        """
        if self.grid and any(any(c for c in r) for r in self.grid):
            return self.constrained_layout()
        if not self.axes.series:
            for s, use_y2 in self.series:
                self.axes.add_series(s, use_y2)
        self.axes.finalize()
        self.axes.tight_layout()
        return self

    def constrained_layout(self) -> "Figure":
        """
        Align subplot cells so their plot areas share a common left edge.

        Equivalent to Matplotlib's ``constrained_layout=True``.  Computes
        the widest Y-tick label across all cells in each column, then sets
        a uniform left padding for all cells in that column so the actual
        data area (to the right of the Y-axis) is flush.

        Returns ``self`` for chaining.

        Example::

            fig = Figure(rows=2, cols=2, width=900, height=600)
            # ... add_axes and add_series ...
            fig.constrained_layout()   # aligns all subplot left edges
        """
        from .utils import _format_tick

        # Finalize all axes first so _y_domain is populated
        for row in self.grid:
            for ax in row:
                if ax is not None and ax.series:
                    ax.finalize()

        # Per-column: compute the max Y-tick label width across all rows
        n_cols = self.cols
        n_rows = self.rows
        col_max_label_px: list[int] = []

        for c in range(n_cols):
            max_label_len = 0
            for r in range(n_rows):
                ax = self.grid[r][c] if r < len(self.grid) and c < len(self.grid[r]) else None
                if ax is None or ax._y_domain is None:
                    continue
                for v in (ax._y_domain[0], ax._y_domain[1]):
                    lbl = _format_tick(v, is_log=(ax.yscale == "log"))
                    max_label_len = max(max_label_len, len(lbl))
            # ~7px per char + 8px gap between label and axis
            col_max_label_px.append(max_label_len * 7 + 8 + 8)

        # Apply uniform padding per column
        for c in range(n_cols):
            uniform_pad = max(col_max_label_px[c], 50)  # never below default 50
            for r in range(n_rows):
                ax = self.grid[r][c] if r < len(self.grid) and c < len(self.grid[r]) else None
                if ax is not None:
                    ax.padding = uniform_pad

        return self

    def add_stat_annotation(
        self,
        x1: Any,
        x2: Any,
        p_value: float = 0.05,
        label: str | None = None,
        style: str = "stars",
        color: str = "#222",
        y_offset: float = 0.0,
    ) -> Figure:
        """
        Add a significance bracket between two groups.

        Draws ``***`` / ``**`` / ``*`` / ``ns`` above the bracket
        based on *p_value*.  Works with both numeric and categorical X axes.

        Args:
            x1:       First group (label or numeric X value).
            x2:       Second group.
            p_value:  Statistical p-value.
            label:    Override the auto-generated significance label.
            style:    ``"stars"`` or ``"numeric"``.
            color:    Bracket and text color.
            y_offset: Extra upward shift in pixels (stack multiple brackets).

        Returns:
            ``self`` for chaining.
        """
        from .stat_annotation import StatAnnotation
        self._stat_annotations = getattr(self, "_stat_annotations", [])
        self._stat_annotations.append(StatAnnotation(
            x1=x1, x2=x2, p_value=p_value,
            label=label, style=style,
            color=color, y_offset=y_offset,
        ))
        return self

    def enable_crosshair(self) -> Figure:
        """
        Enable the synchronized crosshair on the next ``share()`` / ``show()`` call.

        The crosshair draws a vertical line across all GlyphX charts on the
        same HTML page and highlights the nearest data point in each.

        Returns ``self``.
        """
        self._crosshair = True
        return self

    def share(
        self,
        filename: str | None = None,
        title: str | None = None,
    ) -> str:
        """
        Generate a fully self-contained, shareable HTML document.

        The output embeds all JavaScript inline -- no CDN, no server,
        works offline.  Pass ``filename`` to also write it to disk.

        Returns:
            Complete HTML document string.
        """
        from .utils import make_shareable_html
        svg   = self.render_svg()
        label = title or self.title or "GlyphX Chart"
        html  = make_shareable_html(svg, title=label)
        if filename:
            with open(filename, "w", encoding="utf-8") as fh:
                fh.write(html)
        return html

    def plot(self) -> Figure:
        """Shortcut for :meth:`show` respecting ``auto_display``. Returns ``self``."""
        if self.auto_display:
            self.show()
        return self



# ---------------------------------------------------------------------------
# PPTX export helper
# ---------------------------------------------------------------------------

def _save_as_pptx(svg: str, filename: str, title: str | None = None) -> None:
    """
    Save an SVG as a PNG-embedded PowerPoint slide.

    Requires ``python-pptx`` and ``cairosvg``::

        pip install "glyphx[pptx]"

    The SVG is rasterised to PNG at 2x resolution, then inserted as a
    full-slide picture in a blank 16:9 presentation.
    """
    try:
        import cairosvg
    except (ImportError, OSError):
        raise RuntimeError(
            "PPTX export requires cairosvg and the system libcairo library.  "
            "Install with:\n"
            "    pip install \"glyphx[pptx]\"\n"
            "On macOS: brew install cairo"
        )
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        raise RuntimeError(
            "PPTX export requires python-pptx.  Install it with:\n"
            "    pip install \"glyphx[pptx]\""
        )

    import io

    # -- SVG -> PNG at 2x for crisp rendering ------------------------------
    png_bytes = cairosvg.svg2png(bytestring=svg.encode(), scale=2)
    png_stream = io.BytesIO(png_bytes)

    # -- Build presentation ------------------------------------------------
    prs    = Presentation()
    blank  = prs.slide_layouts[6]          # completely blank layout
    slide  = prs.slides.add_slide(blank)

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # -- Optional title text box -------------------------------------------
    top_offset = Inches(0)
    if title:
        txBox = slide.shapes.add_textbox(
            Inches(0.3), Inches(0.1), slide_w - Inches(0.6), Inches(0.55)
        )
        tf = txBox.text_frame
        tf.text = title
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].runs[0].font.size = Pt(22)
        tf.paragraphs[0].runs[0].font.bold = True
        top_offset = Inches(0.65)

    # -- Insert chart PNG --------------------------------------------------
    pic_h = slide_h - top_offset - Inches(0.1)
    pic_w = min(slide_w - Inches(0.4), pic_h * (slide_w / slide_h))
    left  = (slide_w - pic_w) // 2

    slide.shapes.add_picture(png_stream, left, top_offset, pic_w, pic_h)
    prs.save(filename)


# ---------------------------------------------------------------------------
# SubplotGrid
# ---------------------------------------------------------------------------

class SubplotGrid:
    """
    Standalone 2-D grid for laying out existing Figure objects into one page.

    Example::

        sg = SubplotGrid(2, 2)
        sg.add(fig_revenue, 0, 0)
        sg.add(fig_costs,   0, 1)
        html = sg.render()
        open("dashboard.html", "w").write(html)

    Args:
        rows: Number of rows.
        cols: Number of columns.
    """

    def __init__(self, rows: int, cols: int) -> None:
        self.rows = rows
        self.cols = cols
        self.grid: list[list[Figure | None]] = [
            [None] * cols for _ in range(rows)
        ]

    def add(self, figure: Figure, row: int, col: int) -> SubplotGrid:
        """Place a Figure at a grid position. Returns ``self``."""
        if not (0 <= row < self.rows and 0 <= col < self.cols):
            raise IndexError(
                f"Position ({row}, {col}) is out of range for a "
                f"{self.rows}x{self.cols} grid."
            )
        self.grid[row][col] = figure
        return self

    def add_axes(self, row: int, col: int, figure: Figure) -> SubplotGrid:
        """Alias for :meth:`add` kept for backward compatibility."""
        return self.add(figure, row, col)

    def render(self, gap: int = 20) -> str:
        """
        Render all figures into a self-contained HTML page.

        Returns:
            Full HTML document string.
        """
        from .utils import wrap_svg_with_template

        rows_html: list[str] = []
        for r in range(self.rows):
            cells: list[str] = []
            for c in range(self.cols):
                fig = self.grid[r][c]
                svg = fig.render_svg() if fig is not None else ""
                cells.append(f'<div style="margin:{gap}px">{svg}</div>')
            rows_html.append(
                '<div style="display:flex">' + "".join(cells) + "</div>"
            )

        html_body = "<div>" + "".join(rows_html) + "</div>"
        return wrap_svg_with_template(html_body)
